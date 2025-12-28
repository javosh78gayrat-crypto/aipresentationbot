import logging
import sqlite3
import os
import requests
from datetime import datetime, timedelta
from duckduckgo_search import DDGS  # API kalitsiz rasm qidirish uchun

from groq import Groq
from aiogram import Bot, Dispatcher, executor, types
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from aiogram.dispatcher.filters.state import State, StatesGroup
from aiogram.dispatcher import FSMContext
from aiogram.contrib.fsm_storage.memory import MemoryStorage
from aiogram.types import InlineKeyboardMarkup, InlineKeyboardButton

# ================= SOZLAMALAR =================
BOT_TOKEN = "8462890326:AAGh2IJx-ttoAzvvfN4SB7OMdezRVDIt0pY"
GROQ_API_KEY = "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"
CHANNEL_ID = -1003159860704 
ADMIN_ID = 7138813964

SLIDE_PRICE = 5000 
PREMIUM_PRICE = 20000

bot = Bot(token=BOT_TOKEN)
storage = MemoryStorage()
dp = Dispatcher(bot, storage=storage)
groq_client = Groq(api_key=GROQ_API_KEY)

# ================= DATABASE =================
conn = sqlite3.connect("users.db", check_same_thread=False)
cursor = conn.cursor()
cursor.execute("""
CREATE TABLE IF NOT EXISTS users(
    user_id INTEGER PRIMARY KEY,
    balance INT DEFAULT 0,
    premium_until TEXT,
    referals INT DEFAULT 0,
    free_uses INT DEFAULT 0
)""")
conn.commit()

class PresentationStates(StatesGroup):
    waiting_for_topic = State()
    waiting_for_slide_count = State()
    waiting_for_design = State()

# ================= DIZAYN RANGLARI (Aksentlar) =================
DESIGN_COLORS = {
    1: RGBColor(0, 122, 255),   # Apple Blue
    2: RGBColor(255, 45, 85),   # Modern Pink
    3: RGBColor(52, 199, 89),   # Cyber Green
    4: RGBColor(255, 159, 10),  # Bright Orange
    5: RGBColor(175, 82, 222),  # Deep Purple
    6: RGBColor(0, 199, 190),   # Modern Teal
}

# ================= YORDAMCHI FUNKSIYALAR =================
def get_image_url(query):
    try:
        with DDGS() as ddgs:
            results = ddgs.images(query, max_results=1)
            if results: return results[0]['image']
    except: return None
    return None

async def get_ai_content(topic, count):
    # Promptni yanada batafsilroq qildik
    prompt = (f"Prezentatsiya mavzusi: '{topic}'. Jami slaydlar soni: {count}. "
              f"Har bir slayd uchun professional, akademik va batafsil ma'lumot tayyorla. "
              f"Har bir slayd kamida 5 ta kengaytirilgan gapdan iborat bo'lsin. "
              f"Format: 'Title: [Sarlavha]' va 'Content: [Batafsil matn]'. "
              f"Til: O'zbek tili.")
    try:
        chat_completion = groq_client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="llama-3.3-70b-versatile",
            temperature=0.7, # Ijodkorlik va aniqlik balansi
        )
        # ... (qolgan split mantiqi o'zgarishsiz qoladi)
        response = chat_completion.choices[0].message.content
        blocks = response.split("Title:")[1:]
        slides = []
        for b in blocks:
            if "Content:" in b:
                p = b.split("Content:")
                slides.append((p[0].strip(), p[1].strip()))
        return slides[:count]
    except Exception as e:
        logging.error(f"AI Error: {e}")
        return None

def is_premium(user_id):
    cursor.execute("SELECT premium_until FROM users WHERE user_id=?", (user_id,))
    res = cursor.fetchone()
    if not res or res[0] is None: return False
    try:
        return datetime.fromisoformat(res[0]) > datetime.now()
    except: return False

async def check_sub(user_id):
    try:
        m = await bot.get_chat_member(CHANNEL_ID, user_id)
        return m.status in ['member', 'administrator', 'creator']
    except: return False

# ================= START VA MENU =================
@dp.message_handler(commands=['start'])
async def cmd_start(message: types.Message):
    user_id = message.from_user.id
    cursor.execute("INSERT OR IGNORE INTO users (user_id) VALUES (?)", (user_id,))
    
    args = message.get_args()
    if args and args.isdigit() and int(args) != user_id:
        cursor.execute("UPDATE users SET referals = referals + 1 WHERE user_id=?", (args,))
        cursor.execute("SELECT referals FROM users WHERE user_id=?", (args,))
        if cursor.fetchone()[0] >= 15:
            until = (datetime.now() + timedelta(days=30)).isoformat()
            cursor.execute("UPDATE users SET premium_until=?, referals=0 WHERE user_id=?", (until, args))
    conn.commit()

    if not await check_sub(user_id):
        kb = InlineKeyboardMarkup().add(InlineKeyboardButton("✅ Obuna bo'lish", url="https://t.me/ecomind_economy"))
        kb.add(InlineKeyboardButton("🔄 Tekshirish", callback_data="recheck"))
        return await message.answer("⚠️ Botdan foydalanish uchun kanalga obuna bo'ling!", reply_markup=kb)

    main_menu = types.ReplyKeyboardMarkup(resize_keyboard=True, row_width=2)
    main_menu.add("🎤 Prezentatsiya yaratish", "🧾 Tariflar")
    main_menu.add("⭐ Premium", "👥 Referal")
    main_menu.add("📘 Yo‘riqnoma", "💳 Balans")
    await message.answer(f" Assalomu alaykum! Xush kelibsiz 👋 {message.from_user.first_name}!\nPrezentatsiya tayyorlash uchun kerakli bo'limni tanlang.", reply_markup=main_menu)

@dp.message_handler(text="📘 Yo‘riqnoma")
async def help_guide(message: types.Message):
    await message.answer("📖 Bot haqida:\n- Mavzuni yozing\n- Slayd sonini tanlang\n- Dizaynni tanlang\nFayl 1 daqiqada tayyor bo'ladi.\n\nAdmin: @Javohir_fan1")

@dp.message_handler(text="🧾 Tariflar")
async def show_tariffs(message: types.Message):
    await message.answer(f"📊 **Tariflar:**\n\n🎁 Yangi foydalanuvchi: 2 ta bepul dizayn.\n💰 Pullik: 1 slayd = {SLIDE_PRICE} so'm.\n⭐ Premium: {PREMIUM_PRICE} so'm / 1 oy (Cheksiz dizayn va slaydlar).")

@dp.message_handler(text="⭐ Premium")
async def premium_menu(message: types.Message):
    txt = (f"🌟 **PREMIUM TARIFI**\n\n"
           f"✅ 1 oy davomida cheksiz prezentatsiya.\n"
           f"✅ Barcha 7 xil professional dizayn (Glassmorphism & Dark Mode).\n"
           f"✅ Slaydlar sonini tanlash (30 tagacha).\n\n"
           f"💰 Narxi: {PREMIUM_PRICE} so'm\n"
           f"💳 Karta: `4073420043005051`\n"
           f"👤 Ism: Javohir G'ayrataliyev\n\n"
           f"To'lov chekini yuboring!")
    await message.answer(txt, parse_mode="Markdown")

@dp.message_handler(text="👥 Referal")
async def referal_menu(message: types.Message):
    bot_name = (await bot.get_me()).username
    link = f"https://t.me/{bot_name}?start={message.from_user.id}"
    cursor.execute("SELECT referals FROM users WHERE user_id=?", (message.from_user.id,))
    count = cursor.fetchone()[0]
    await message.answer(f"🔗 Referal linkingiz:\n{link}\n\nTakliflar: {count}/15\n🎁 15 do'st = 1 oy bepul Premium!")

@dp.message_handler(text="💳 Balans")
async def show_balance(message: types.Message):
    cursor.execute("SELECT balance, free_uses FROM users WHERE user_id=?", (message.from_user.id,))
    res = cursor.fetchone()
    await message.answer(f"💳 Balansingiz: {res[0]} so'm\n🎁 Tekin imkoniyat: {2-res[1] if res[1]<2 else 0} ta")

# ================= PPT TAYYORLASH FLOW =================
@dp.message_handler(text="🎤 Prezentatsiya yaratish")
async def start_ppt(message: types.Message):
    if not await check_sub(message.from_user.id): return
    await message.answer("📝 Prezentatsiya mavzusini kiriting:")
    await PresentationStates.waiting_for_topic.set()

@dp.message_handler(state=PresentationStates.waiting_for_topic)
async def get_topic(message: types.Message, state: FSMContext):
    await state.update_data(topic=message.text)
    user_id = message.from_user.id
    if is_premium(user_id):
        await message.answer("🔢 Slaydlar sonini kiriting (1-30):")
        await PresentationStates.waiting_for_slide_count.set()
    else:
        cursor.execute("SELECT balance, free_uses FROM users WHERE user_id=?", (user_id,))
        bal, free = cursor.fetchone()
        if free < 2 or bal >= (SLIDE_PRICE * 10):
            await state.update_data(slide_count=10)
            await ask_design(message, user_id)
        else:
            await message.answer("❌ Limit tugagan yoki mablag' yetarli emas. Premium sotib oling!")
            await state.finish()

async def ask_design(message, user_id):
    kb = InlineKeyboardMarkup(row_width=2)
    limit = 8 if is_premium(user_id) else 3 
    for i in range(1, limit):
        btn_text = f"🎨 Dizayn {i}" if i < 7 else "💎 Glassmorphism"
        kb.insert(InlineKeyboardButton(btn_text, callback_data=f"ds_{i}"))
    
    if not is_premium(user_id):
        kb.add(InlineKeyboardButton("⭐ Barcha dizaynlar (Premium)", callback_data="premium_info"))
        
    await message.answer("🎨 Dizaynni tanlang:", reply_markup=kb)
    await PresentationStates.waiting_for_design.set()

@dp.message_handler(state=PresentationStates.waiting_for_slide_count)
async def get_count(message: types.Message, state: FSMContext):
    if message.text.isdigit() and 1 <= int(message.text) <= 30:
        await state.update_data(slide_count=int(message.text))
        await ask_design(message, message.from_user.id)
    else:
        await message.answer("⚠️ 1-30 oralig'ida raqam kiriting!")

# ! FINAL MUKAMMAL DIZAYN GENERATSIYASI !
@dp.callback_query_handler(lambda c: c.data.startswith('ds_'), state=PresentationStates.waiting_for_design)
async def finalize_ppt(call: types.CallbackQuery, state: FSMContext):
    design_id = int(call.data.split('_')[1])
    data = await state.get_data()
    user_id = call.from_user.id
    
    msg = await call.message.edit_text("✨ Zamonaviy dizayn va rasmlar tayyorlanmoqda...")
    slides_data = await get_ai_content(data['topic'], data['slide_count'])
    
    if not slides_data:
        await msg.edit_text("❌ Xatolik! AI ma'lumot bera olmadi.")
        return await state.finish()

    prs = Presentation()
    # Dizayn ranglari (7-dizayn premium dark mode)
    accent_rgb = DESIGN_COLORS.get(design_id if design_id < 7 else 1, RGBColor(0, 122, 255))
    bg_rgb = RGBColor(18, 18, 18)  # Modern Apple Dark #121212
    text_main = RGBColor(240, 240, 240) # Off-white

    slide_num = 1
    for title_text, content_text in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # --- DARK MODE FON (HAMMA UCHUN) ---
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_rgb

        # 1. YUQORI GEOMETRIK CHIZIQ
        top_line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = accent_rgb
        top_line.line.fill.background()

        # 2. RASM QO'SHISH (DUCKDUCKGO) - O'NG TOMONGA
        img_url = get_image_url(f"{data['topic']} {title_text}")
        content_width = Inches(8.8)
        if img_url:
            try:
                img_res = requests.get(img_url, timeout=5)
                temp_img = f"img_{user_id}_{slide_num}.jpg"
                with open(temp_img, "wb") as f: f.write(img_res.content)
                # Rasmni zamonaviy ramkada o'ngga qo'yish
                slide.shapes.add_picture(temp_img, Inches(6.2), Inches(1.3), Inches(3.4), Inches(4.8))
                content_width = Inches(5.5)
                os.remove(temp_img)
            except: pass

        # 3. SARLAVHA (UPPERCASE & BOLD)
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(0.8))
        tf_t = title_box.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text.upper()
        p_t.font.bold = True
        p_t.font.size = Pt(26)
        p_t.font.color.rgb = accent_rgb

        # 4. SARLAVHA OSTIDAGI QISQA CHIZIQ
        und_line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.1), Inches(1.2), Inches(0.04))
        und_line.fill.solid()
        und_line.fill.fore_color.rgb = accent_rgb
        und_line.line.fill.background()

        # 5. ASOSIY MATN (MODERN MARKERLAR ▹)
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), content_width, Inches(5))
        tf_c = content_box.text_frame
        tf_c.word_wrap = True
        for line in content_text.split('.'):
            if len(line.strip()) > 5:
                p_c = tf_c.add_paragraph()
                p_c.text = "▹ " + line.strip()
                p_c.font.size = Pt(16)
                p_c.font.color.rgb = text_main
                p_c.space_before = Pt(12)

        # 6. FOOTER (BRAND & PAGE)
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(5), Inches(0.3))
        p_f = footer_box.text_frame.paragraphs[0]
        p_f.text = f"© AI Presentation Bot | Slayd {slide_num}"
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = RGBColor(120, 120, 120)
        slide_num += 1

    file_name = f"modern_ppt_{user_id}.pptx"
    prs.save(file_name)
    await bot.send_document(user_id, open(file_name, 'rb'), caption=f"✅ {data['topic']} uchun zamonaviy prezentatsiya tayyor!")
    os.remove(file_name)
    await msg.delete()

    if not is_premium(user_id):
        cursor.execute("SELECT free_uses FROM users WHERE user_id=?", (user_id,))
        if cursor.fetchone()[0] < 2:
            cursor.execute("UPDATE users SET free_uses = free_uses + 1 WHERE user_id=?", (user_id,))
        else:
            cursor.execute("UPDATE users SET balance = balance - ? WHERE user_id=?", (data['slide_count']*SLIDE_PRICE, user_id))
    conn.commit()
    await state.finish()

# ================= ADMIN VA BOSHQALAR =================
@dp.callback_query_handler(lambda c: c.data == "premium_info", state="*")
async def buy_premium_call(call: types.CallbackQuery):
    await premium_menu(call.message)

@dp.message_handler(content_types=['photo'])
async def handle_payment(message: types.Message):
    kb = InlineKeyboardMarkup(row_width=1)
    kb.add(
        InlineKeyboardButton("✅ Premium berish", callback_data=f"adm_ok_p_{message.from_user.id}"),
        InlineKeyboardButton("✅ 20 ming balans", callback_data=f"adm_ok_b_{message.from_user.id}"),
        InlineKeyboardButton("❌ Rad etish", callback_data=f"adm_no_x_{message.from_user.id}")
    )
    await bot.send_photo(ADMIN_ID, message.photo[-1].file_id, 
                         caption=f"Foydalanuvchi: {message.from_user.full_name}\nID: {message.from_user.id}", 
                         reply_markup=kb)
    await message.answer("📸 Chek yuborildi. Admin tasdiqlashini kuting.")

@dp.callback_query_handler(lambda c: c.data.startswith('adm_'))
async def admin_callback(call: types.CallbackQuery):
    parts = call.data.split('_')
    status, type_, u_id = parts[1], parts[2], parts[3]
    if status == 'ok':
        if type_ == 'p':
            until = (datetime.now() + timedelta(days=30)).isoformat()
            cursor.execute("UPDATE users SET premium_until=? WHERE user_id=?", (until, u_id))
            await bot.send_message(u_id, "💎 Premium faollashtirildi (30 kun)!")
        else:
            cursor.execute("UPDATE users SET balance = balance + 20000 WHERE user_id=?", (u_id,))
            await bot.send_message(u_id, "💰 Balansingiz 20,000 so'mga to'ldirildi!")
    else:
        await bot.send_message(u_id, "❌ To'lov cheki admin tomonidan rad etildi.")
    conn.commit()
    await call.message.edit_caption("Bajarildi ✅")

@dp.callback_query_handler(lambda c: c.data == "recheck")
async def recheck(call: types.CallbackQuery):
    if await check_sub(call.from_user.id):
        await call.message.delete()
        await cmd_start(call.message)
    else:
        await call.answer("❌ Kanalga obuna bo'lmagansiz!", show_alert=True)

if __name__ == "__main__":
    executor.start_polling(dp, skip_updates=True)
